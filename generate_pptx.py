from pptx import Presentation

# Path to your outline file
outline_path = r'd:\SCM PROJECT\PerishAI_Presentation.txt'

# Create a new presentation
prs = Presentation()

with open(outline_path, 'r', encoding='utf-8') as f:
    lines = f.readlines()

slide = None
for line in lines:
    line = line.strip()
    if line.startswith('Slide'):
        # Start a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
    elif line.startswith('---') or not line:
        continue
    elif slide and ':' in line:
        # Slide title
        title = line.split(':', 1)[1].strip()
        slide.shapes.title.text = title
    elif slide:
        # Slide content
        content = slide.placeholders[1]
        content.text += line + '\n'

# Save the presentation
prs.save(r'd:\SCM PROJECT\PerishAI_Presentation.pptx')
print('PowerPoint file created: d:\SCM PROJECT\PerishAI_Presentation.pptx')
